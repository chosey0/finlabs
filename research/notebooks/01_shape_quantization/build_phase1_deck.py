"""Build a comprehensive PPT deck for the 01_shape_quantization Phase 1 results.

Run with:
    uv run --with python-pptx python research/notebooks/01_shape_quantization/build_phase1_deck.py

Content is sourced directly from results/*.md and summaries/SUMMARY.md.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).resolve().parent
P1A = BASE / "runs/phase_1a/price_shape_NASDAQ_1m_k12/figures"
P1B_FIX = BASE / "runs/phase_1b/shape_token_range_bucket_NASDAQ_1m_k12/figures"
P1B_RND = (
    BASE / "runs/phase_1b/shape_token_range_bucket_NASDAQ_1m_k12_random_00/figures"
)
P1B_ABL = BASE / "runs/phase_1b/shape_range_NASDAQ_1m_k12/figures"
OUT = BASE / "Phase1_Shape_Quantization_Results.pptx"

# ---- palette -------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x2A, 0x44)
NAVY_2 = RGBColor(0x2C, 0x3A, 0x5E)
BLUE = RGBColor(0x4C, 0x7D, 0xF0)
UP = RGBColor(0xFD, 0x79, 0x79)  # korean candle up
DOWN = RGBColor(0x8C, 0xA9, 0xFF)  # korean candle down
GREEN = RGBColor(0x2E, 0xA8, 0x6B)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
GREY = RGBColor(0x6B, 0x74, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARKTEXT = RGBColor(0x22, 0x28, 0x36)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- helpers -------------------------------------------------------------
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, left, t, w, h, fill=None, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE

    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def _text(
    slide,
    left,
    t,
    w,
    h,
    runs,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    space_after=6,
    line_spacing=1.05,
):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(left, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for txt, size, color, bold in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def header(slide, kicker, title, dark=False):
    _bg(slide, NAVY if dark else WHITE)
    # accent bar
    _box(slide, Inches(0.55), Inches(0.55), Inches(0.12), Inches(0.95), fill=BLUE)
    tcol = WHITE if dark else NAVY
    kcol = DOWN if dark else BLUE
    _text(
        slide,
        Inches(0.85),
        Inches(0.5),
        Inches(11.8),
        Inches(0.45),
        [[(kicker, 13, kcol, True)]],
    )
    _text(
        slide,
        Inches(0.85),
        Inches(0.85),
        Inches(11.8),
        Inches(0.8),
        [[(title, 28, tcol, True)]],
    )


def footer(slide, idx, dark=False):
    col = DOWN if dark else GREY
    _text(
        slide,
        Inches(0.85),
        Inches(7.02),
        Inches(8),
        Inches(0.3),
        [
            [
                (
                    "FinLabs Research · Phase 1 Shape Quantization · NASDAQ 1m",
                    9,
                    col,
                    False,
                )
            ]
        ],
    )
    _text(
        slide,
        Inches(11.8),
        Inches(7.02),
        Inches(1.2),
        Inches(0.3),
        [[(str(idx), 10, col, True)]],
        align=PP_ALIGN.RIGHT,
    )


def table(
    slide,
    left,
    t,
    w,
    rows,
    col_w=None,
    header_fill=NAVY,
    font=11,
    first_col_left=True,
    zebra=True,
    highlight_rows=None,
):
    nrows = len(rows)
    ncols = len(rows[0])
    h = Inches(0.34 * nrows)
    gt = slide.shapes.add_table(nrows, ncols, left, t, w, h).table
    if col_w:
        total = sum(col_w)
        for j, cw in enumerate(col_w):
            gt.columns[j].width = Emu(int(w * cw / total))
    highlight_rows = highlight_rows or set()
    for i, row in enumerate(rows):
        gt.rows[i].height = Inches(0.34)
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.06)
            c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02)
            c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (
                PP_ALIGN.LEFT if (j == 0 and first_col_left) else PP_ALIGN.CENTER
            )
            r = p.add_run()
            r.text = str(val)
            r.font.name = "Calibri"
            if i == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = header_fill
                r.font.size = Pt(font)
                r.font.bold = True
                r.font.color.rgb = WHITE
            else:
                if i in highlight_rows:
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(0xE2, 0xEC, 0xFF)
                elif zebra and i % 2 == 0:
                    c.fill.solid()
                    c.fill.fore_color.rgb = LIGHT
                else:
                    c.fill.solid()
                    c.fill.fore_color.rgb = WHITE
                r.font.size = Pt(font)
                r.font.bold = (j == 0 and first_col_left) or (i in highlight_rows)
                r.font.color.rgb = DARKTEXT
    return gt


def picture(slide, path, left, t, w=None, h=None, caption=None):
    pic = slide.shapes.add_picture(str(path), left, t, width=w, height=h)
    if caption:
        _text(
            slide,
            left,
            pic.top + pic.height + Inches(0.04),
            pic.width,
            Inches(0.3),
            [[(caption, 9.5, GREY, False)]],
            align=PP_ALIGN.CENTER,
        )
    return pic


def bullets(slide, left, t, w, h, items, size=14, gap=8):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            txt, lvl, color, bold = it
        else:
            txt, lvl, color, bold = it, 0, DARKTEXT, False
        prefix = "    " * lvl + ("▪  " if lvl == 0 else "–  ")
        paras.append(
            [(prefix, size, BLUE if lvl == 0 else GREY, True), (txt, size, color, bold)]
        )
    _text(slide, left, t, w, h, paras, space_after=gap, line_spacing=1.05)


def card(slide, left, t, w, h, label, value, vcolor=BLUE, sub=None):
    _box(
        slide,
        left,
        t,
        w,
        h,
        fill=LIGHT,
        line=RGBColor(0xDD, 0xE3, 0xEE),
        line_w=Pt(0.75),
    )
    _text(
        slide,
        left,
        t + Inches(0.12),
        w,
        Inches(0.35),
        [[(label, 11, GREY, True)]],
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        left,
        t + Inches(0.42),
        w,
        Inches(0.7),
        [[(value, 30, vcolor, True)]],
        align=PP_ALIGN.CENTER,
    )
    if sub:
        _text(
            slide,
            left,
            t + h - Inches(0.42),
            w,
            Inches(0.35),
            [[(sub, 10, GREY, False)]],
            align=PP_ALIGN.CENTER,
        )


n = 0


def new(dark=False):
    global n
    n += 1
    s = prs.slides.add_slide(BLANK)
    return s


# ========================================================================
# 1. TITLE
# ========================================================================
s = new(dark=True)
_bg(s, NAVY)
_box(s, 0, Inches(6.7), SW, Inches(0.8), fill=NAVY_2)
_box(s, Inches(0.9), Inches(2.35), Inches(0.16), Inches(2.0), fill=BLUE)
_text(
    s,
    Inches(1.25),
    Inches(2.25),
    Inches(11),
    Inches(0.5),
    [[("FINLABS RESEARCH · 01_SHAPE_QUANTIZATION", 15, DOWN, True)]],
)
_text(
    s,
    Inches(1.25),
    Inches(2.8),
    Inches(11.2),
    Inches(1.6),
    [
        [("Shape Quantization (Phase 1)", 44, WHITE, True)],
        [
            (
                "캔들 price-shape의 discrete tokenization 종합 결과",
                24,
                RGBColor(0xC9, 0xD4, 0xEC),
                False,
            )
        ],
    ],
    space_after=10,
)
_text(
    s,
    Inches(1.25),
    Inches(5.0),
    Inches(11),
    Inches(1.0),
    [
        [
            (
                "VQ-VAE Candlestick Tokenizer  ·  NASDAQ 1-minute  ·  Codebook K=12",
                16,
                DOWN,
                False,
            )
        ],
        [
            (
                "representation = (shape_token, range_bucket)   |   primary filter: volume ≥ 2",
                13,
                GREY,
                False,
            )
        ],
    ],
    space_after=6,
)
_text(
    s,
    Inches(1.25),
    Inches(6.85),
    Inches(11),
    Inches(0.5),
    [[("2026-06-05", 11, RGBColor(0xC9, 0xD4, 0xEC), False)]],
)

# ========================================================================
# 2. RESEARCH FRAME / STATUS
# ========================================================================
s = new()
header(s, "RESEARCH FRAME", "3단계 안전장치와 현재 위치")
bullets(
    s,
    Inches(0.85),
    Inches(1.75),
    Inches(5.7),
    Inches(3.2),
    [
        (
            "시장 데이터를 price prediction target이 아니라 학습 가능한",
            0,
            DARKTEXT,
            False,
        ),
        ("market representation으로 재구성", 1, NAVY, True),
        (
            "결론을 처음부터 'market state'로 단정하지 않고 3단계로 분리",
            0,
            DARKTEXT,
            False,
        ),
        ("Phase 1 산출물은 shape token (state 아님)", 1, GREY, False),
        ("Phase 2 전이 구조 확인 시 state candidate", 1, GREY, False),
        ("Phase 3 미래 dynamics 검증 후에야 market state", 1, GREY, False),
    ],
    size=14,
    gap=9,
)
table(
    s,
    Inches(6.95),
    Inches(1.85),
    Inches(5.55),
    [
        ["Phase", "내용", "상태"],
        ["1A Price-shape only", "4D shape 양자화", "완료"],
        ["1B Shape+Range bucket", "반복 split 검증", "통과"],
        ["2 Sequential Dynamics", "token 전이 구조", "계획 수립"],
        ["3 Market State", "미래 dynamics", "future"],
    ],
    col_w=[2.4, 2.2, 1.1],
    font=12,
    highlight_rows={1, 2},
)
_box(
    s,
    Inches(6.95),
    Inches(5.5),
    Inches(5.55),
    Inches(1.05),
    fill=RGBColor(0xE7, 0xF3, 0xEC),
    line=GREEN,
    line_w=Pt(1),
)
_text(
    s,
    Inches(7.15),
    Inches(5.62),
    Inches(5.2),
    Inches(0.85),
    [
        [
            ("현재 위치: ", 13, GREEN, True),
            ("Phase 1B 통과 → Phase 2 진입 기준 충족", 13, NAVY, True),
        ],
        [("본 자료는 Phase 1 (1A+1B) 전체 결과의 종합입니다.", 11, GREY, False)],
    ],
    space_after=3,
)
footer(s, n)

# ========================================================================
# 3. QUESTION & NON-CLAIMS
# ========================================================================
s = new(dark=True)
header(s, "SCOPE", "Phase 1 질문과 의도적 비주장", dark=True)
_box(s, Inches(0.85), Inches(1.8), Inches(11.6), Inches(1.15), fill=NAVY_2)
_text(
    s,
    Inches(1.15),
    Inches(1.95),
    Inches(11),
    Inches(0.9),
    [
        [("Q. 여러 NASDAQ symbol에서 반복되는 candle price-shape를", 18, WHITE, True)],
        [("    같은 discrete token으로 묶을 수 있는가?", 18, DOWN, True)],
    ],
    space_after=4,
)
_text(
    s,
    Inches(0.85),
    Inches(3.25),
    Inches(5.6),
    Inches(0.4),
    [[("✓ 이 단계에서 다루는 것", 15, GREEN, True)]],
)
bullets(
    s,
    Inches(0.85),
    Inches(3.75),
    Inches(5.6),
    Inches(3),
    [
        (
            "상대적 candle price-shape vocabulary 학습",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
        (
            "held-out symbol에서의 token 분포 안정성",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
        ("volatility context의 분리 가능성", 0, RGBColor(0xD9, 0xE2, 0xF2), False),
    ],
    size=14,
    gap=10,
)
_text(
    s,
    Inches(6.8),
    Inches(3.25),
    Inches(5.6),
    Inches(0.4),
    [[("✗ 아직 주장하지 않는 것", 15, UP, True)]],
)
bullets(
    s,
    Inches(6.8),
    Inches(3.75),
    Inches(5.6),
    Inches(3),
    [
        ("market state discovery", 0, RGBColor(0xD9, 0xE2, 0xF2), False),
        ("future return / volatility 예측", 0, RGBColor(0xD9, 0xE2, 0xF2), False),
        ("trading signal · 매수/매도 조건", 0, RGBColor(0xD9, 0xE2, 0xF2), False),
        ("VQ-VAE가 KMeans보다 우월하다는 결론", 0, RGBColor(0xD9, 0xE2, 0xF2), False),
    ],
    size=14,
    gap=10,
)
footer(s, n, dark=True)

# ========================================================================
# 4. EXPERIMENT SETUP
# ========================================================================
s = new()
header(s, "SETUP", "공통 실험 설정")
table(
    s,
    Inches(0.85),
    Inches(1.8),
    Inches(5.7),
    [
        ["Item", "Value"],
        ["Market / Interval", "NASDAQ / 1m"],
        ["Max candles / symbol", "12,000"],
        ["Min candles / symbol", "500"],
        ["Codebook size K", "12"],
        ["Epochs / Batch / LR", "25 / 256 / 0.001"],
        ["Primary candle filter", "volume ≥ 2"],
    ],
    col_w=[2.3, 2.4],
    font=12,
)
_text(
    s,
    Inches(0.85),
    Inches(4.65),
    Inches(5.7),
    Inches(0.4),
    [[("4D price-shape feature", 14, BLUE, True)]],
)
bullets(
    s,
    Inches(0.85),
    Inches(5.1),
    Inches(5.7),
    Inches(1.8),
    [
        ("signed_body_ratio  ·  upper_ratio", 0, DARKTEXT, False),
        ("lower_ratio  ·  body_center_location", 0, DARKTEXT, False),
        (
            "제외: volume, 절대 price level, range scale, future return, label",
            0,
            GREY,
            False,
        ),
    ],
    size=12,
    gap=6,
)
_text(
    s,
    Inches(6.95),
    Inches(1.8),
    Inches(5.6),
    Inches(0.4),
    [[("Held-out symbol split (fixed)", 14, BLUE, True)]],
)
table(
    s,
    Inches(6.95),
    Inches(2.25),
    Inches(5.55),
    [
        ["Split", "Symbols", "Candles"],
        ["train", "AAPL MSFT NVDA TSLA AMZN META GOOGL AMD INTC RKLB AVGO", "132,000"],
        ["val", "NFLX PLTR", "24,000"],
        ["test", "MU QCOM MRVL", "36,000"],
    ],
    col_w=[0.8, 3.3, 1.2],
    font=10.5,
)
_box(
    s,
    Inches(6.95),
    Inches(4.55),
    Inches(5.55),
    Inches(2.0),
    fill=RGBColor(0xFF, 0xF3, 0xE9),
    line=UP,
    line_w=Pt(1),
)
_text(
    s,
    Inches(7.15),
    Inches(4.68),
    Inches(5.2),
    Inches(1.8),
    [
        [("Leakage 방지 규칙", 13, UP, True)],
        [("• VQ-VAE는 train candles로만 학습", 12, DARKTEXT, False)],
        [("• RangeBucketizer quantile도 train으로만 fit", 12, DARKTEXT, False)],
        [("• val/test는 train 통계로만 transform", 12, DARKTEXT, False)],
        [("• symbol 경계를 넘는 transition 생성 금지", 12, DARKTEXT, False)],
    ],
    space_after=5,
)
footer(s, n)

# ========================================================================
# 5. PHASE 1A VQ-VAE SUMMARY
# ========================================================================
s = new()
header(s, "PHASE 1A", "Price-shape only 양자화 — VQ-VAE 결과")
table(
    s,
    Inches(0.85),
    Inches(1.8),
    Inches(7.1),
    [
        ["Metric", "K=8", "K=12", "K=16"],
        ["Used tokens (train)", "8 / 8", "12 / 12", "13 / 16"],
        ["Dead tokens", "0", "0", "3"],
        ["Train entropy", "2.947", "3.469", "3.581"],
        ["Test–train L1 (drift)", "0.074", "0.108", "0.108"],
        ["Semantic consistency", "0.259", "0.196", "0.190"],
        ["Reconstruction MSE", "0.0199", "0.0127", "0.0118"],
    ],
    col_w=[2.4, 1.3, 1.3, 1.3],
    font=12,
    highlight_rows={2},
)
_text(s, Inches(0.85), Inches(0.0), Inches(1), Inches(0.1), [[("", 1, WHITE, False)]])
card(
    s,
    Inches(8.25),
    Inches(1.85),
    Inches(4.25),
    Inches(1.55),
    "채택 설정",
    "K = 12",
    vcolor=BLUE,
    sub="no dead token · interpretable vocabulary",
)
bullets(
    s,
    Inches(8.25),
    Inches(3.65),
    Inches(4.25),
    Inches(3),
    [
        ("K=8: 안정적이나 vocabulary가 coarse", 0, DARKTEXT, False),
        ("K=12: dead token 0, 의미·복원 개선", 0, NAVY, True),
        ("K=16: 복원은 더 좋지만 dead token 3개", 0, DARKTEXT, False),
    ],
    size=12.5,
    gap=9,
)
bullets(
    s,
    Inches(0.85),
    Inches(4.55),
    Inches(7.1),
    Inches(2),
    [
        (
            "split drift는 모든 K에서 낮음 (K=12 test–train L1 = 0.108)",
            0,
            DARKTEXT,
            False,
        ),
        (
            "결론: 반복되는 candle shape를 discrete shape token 후보로 묶을 수 있다",
            0,
            NAVY,
            True,
        ),
    ],
    size=13,
    gap=8,
)
footer(s, n)

# ========================================================================
# 6. PHASE 1A KMEANS BASELINE
# ========================================================================
s = new()
header(s, "PHASE 1A", "KMeans baseline 비교")
table(
    s,
    Inches(0.85),
    Inches(1.85),
    Inches(7.0),
    [
        ["Metric (K=12)", "VQ-VAE", "KMeans"],
        ["Used clusters", "12 / 12", "12 / 12"],
        ["Train entropy", "3.469", "3.539"],
        ["Val–train L1", "0.089", "0.077"],
        ["Test–train L1", "0.108", "0.096"],
        ["Semantic consistency", "0.196", "0.210"],
    ],
    col_w=[2.6, 1.5, 1.5],
    font=12.5,
)
picture(
    s,
    P1A / "08_vqvae_vs_kmeans_histogram.png",
    Inches(8.1),
    Inches(1.85),
    w=Inches(4.5),
    caption="K=12 · VQ-VAE vs KMeans token ratio",
)
_box(
    s,
    Inches(0.85),
    Inches(4.95),
    Inches(7.0),
    Inches(1.6),
    fill=RGBColor(0xFF, 0xF3, 0xE9),
    line=UP,
    line_w=Pt(1),
)
_text(
    s,
    Inches(1.05),
    Inches(5.08),
    Inches(6.6),
    Inches(1.4),
    [
        [("4D handcrafted feature에서는 KMeans가 강한 baseline", 13, UP, True)],
        [
            (
                "현 수치만으로 VQ-VAE 우위를 주장할 수 없음. VQ-VAE 유지 근거는",
                12,
                DARKTEXT,
                False,
            )
        ],
        [
            (
                "단순 clustering 성능이 아니라 sequential modeling 확장성에 있음.",
                12,
                DARKTEXT,
                False,
            )
        ],
    ],
    space_after=4,
)
footer(s, n)

# ========================================================================
# 7. PHASE 1A VISUAL
# ========================================================================
s = new(dark=True)
header(s, "PHASE 1A · VISUAL EVIDENCE", "학습된 shape vocabulary (K=12)", dark=True)
picture(
    s,
    P1A / "06_prototype_candles.png",
    Inches(0.85),
    Inches(1.85),
    w=Inches(5.7),
    caption="Prototype candles — 12개 shape token 원형",
)
picture(
    s,
    P1A / "05_mean_feature_heatmap.png",
    Inches(6.9),
    Inches(1.85),
    w=Inches(5.7),
    caption="Mean feature heatmap — token별 4D feature 평균",
)
footer(s, n, dark=True)

# ========================================================================
# 8. PHASE 1B DESIGN + ABLATION
# ========================================================================
s = new()
header(s, "PHASE 1B", "설계 — shape token과 range bucket 분리")
_box(
    s,
    Inches(0.85),
    Inches(1.8),
    Inches(11.6),
    Inches(0.95),
    fill=RGBColor(0xE2, 0xEC, 0xFF),
    line=BLUE,
    line_w=Pt(1),
)
_text(
    s,
    Inches(1.05),
    Inches(1.92),
    Inches(11.2),
    Inches(0.8),
    [
        [
            ("final representation = ", 16, NAVY, True),
            ("(shape_token, range_bucket)", 16, BLUE, True),
            (
                "    — shape는 4D VQ-VAE token, volatility는 quantile bucket으로 분리",
                13,
                GREY,
                False,
            ),
        ]
    ],
)
_text(
    s,
    Inches(0.85),
    Inches(3.0),
    Inches(5.6),
    Inches(0.4),
    [[("✗ Ablation: range를 encoder input에 직접 투입", 14, UP, True)]],
)
table(
    s,
    Inches(0.85),
    Inches(3.5),
    Inches(5.6),
    [
        ["Metric (5D input)", "VQ-VAE", "KMeans"],
        ["Val–train L1", "0.135", "0.130"],
        ["Test–train L1", "0.627", "0.658"],
        ["Semantic consistency", "0.533", "0.517"],
    ],
    col_w=[2.4, 1.3, 1.3],
    font=12,
    highlight_rows={2},
)
_text(
    s,
    Inches(0.85),
    Inches(5.25),
    Inches(5.6),
    Inches(1.3),
    [
        [
            (
                "range_scale_z를 섞으면 token이 shape와 volatility를 동시에 담아",
                12,
                DARKTEXT,
                False,
            )
        ],
        [("test drift가 0.627로 급증 → shape vocabulary 의미가 흐려짐.", 12, UP, True)],
    ],
    space_after=4,
)
_text(
    s,
    Inches(6.8),
    Inches(3.0),
    Inches(5.7),
    Inches(0.4),
    [[("✓ 채택: shape_token + 별도 range_bucket", 14, GREEN, True)]],
)
bullets(
    s,
    Inches(6.8),
    Inches(3.5),
    Inches(5.7),
    Inches(3),
    [
        (
            "range_pct = (high−low)/ref ;  log1p 변환 후 quantile bucket",
            0,
            DARKTEXT,
            False,
        ),
        ("buckets: very_low/low/normal/high/very_high/extreme", 0, GREY, False),
        ("quantile 경계 20/40/60/80/95% (train으로만 fit)", 0, GREY, False),
        ("의미 분리 → shape drift와 volatility drift를 분해 관찰 가능", 0, NAVY, True),
    ],
    size=12.5,
    gap=8,
)
footer(s, n)

# ========================================================================
# 9. PHASE 1B FIXED SPLIT RESULTS
# ========================================================================
s = new()
header(s, "PHASE 1B", "Fixed split 결과 — shape는 안정, range/pair는 drift")
table(
    s,
    Inches(0.85),
    Inches(1.85),
    Inches(11.6),
    [
        ["대상", "Used / Total", "Test–train L1", "해석"],
        [
            "shape_token",
            "12 / 12 (dead 0)",
            "0.108",
            "Phase 1A와 동일 수준 — vocabulary 유지",
        ],
        [
            "range_bucket",
            "6 / 6",
            "0.674",
            "held-out symbol의 volatility profile 차이 반영",
        ],
        [
            "shape × range pair",
            "67 / 72",
            "0.712",
            "drift 대부분이 range_bucket에서 유입",
        ],
    ],
    col_w=[1.8, 1.8, 1.6, 4.6],
    font=12.5,
    highlight_rows={1},
)
bullets(
    s,
    Inches(0.85),
    Inches(3.85),
    Inches(11.6),
    Inches(1.6),
    [
        (
            "shape token은 12개 모두 사용, dead token 없음 → range 분리해도 shape vocabulary 보존",
            0,
            NAVY,
            True,
        ),
        (
            "pair drift는 'shape drift'가 아니라 'symbol-level range profile drift'로 해석",
            0,
            DARKTEXT,
            False,
        ),
    ],
    size=13,
    gap=8,
)
_text(
    s,
    Inches(0.85),
    Inches(5.0),
    Inches(5.5),
    Inches(0.4),
    [[("VQ-VAE vs KMeans (shape test–train L1)", 12, GREY, True)]],
)
table(
    s,
    Inches(0.85),
    Inches(5.45),
    Inches(6.2),
    [
        ["Model", "shape val", "shape test", "pair test"],
        ["VQ-VAE", "0.089", "0.108", "0.712"],
        ["KMeans", "0.077", "0.096", "0.701"],
    ],
    col_w=[1.6, 1.4, 1.4, 1.4],
    font=11.5,
)
footer(s, n)

# ========================================================================
# 10. PHASE 1B FIXED VISUAL
# ========================================================================
s = new(dark=True)
header(
    s,
    "PHASE 1B · VISUAL EVIDENCE",
    "Shape는 symbol 간 안정, range는 symbol별로 분화",
    dark=True,
)
picture(
    s,
    P1B_FIX / "05_per_symbol_shape_token_heatmap.png",
    Inches(0.85),
    Inches(1.85),
    w=Inches(5.7),
    caption="Per-symbol shape token — symbol 간 분포 유사 (안정)",
)
picture(
    s,
    P1B_FIX / "04_per_symbol_range_bucket_heatmap.png",
    Inches(6.9),
    Inches(1.85),
    w=Inches(5.7),
    caption="Per-symbol range bucket — symbol별 volatility 분화",
)
footer(s, n, dark=True)

# ========================================================================
# 11. PHASE 1B REPEATED RANDOM (initial 5)
# ========================================================================
s = new()
header(s, "PHASE 1B", "반복 random split 검증 (초기 5회)")
table(
    s,
    Inches(0.85),
    Inches(1.8),
    Inches(7.7),
    [
        ["Metric", "Mean", "Std", "Min", "Max"],
        ["shape val–train L1", "0.155", "0.060", "0.079", "0.234"],
        ["shape test–train L1", "0.107", "0.050", "0.050", "0.148"],
        ["range val–train L1", "0.475", "0.247", "0.135", "0.729"],
        ["pair val–train L1", "0.512", "0.238", "0.176", "0.773"],
        ["KMeans shape test L1", "0.103", "0.050", "0.050", "0.160"],
    ],
    col_w=[2.6, 1.1, 1.1, 1.1, 1.1],
    font=12,
    highlight_rows={2},
)
card(
    s,
    Inches(8.85),
    Inches(1.85),
    Inches(3.65),
    Inches(1.7),
    "shape test drift",
    "0.107",
    vcolor=GREEN,
    sub="fixed split 0.108과 동일 수준",
)
bullets(
    s,
    Inches(8.85),
    Inches(3.8),
    Inches(3.65),
    Inches(3),
    [
        ("train/test symbol을 무작위로 바꿔도", 0, DARKTEXT, False),
        ("shape token 분포는 흔들리지 않음", 1, NAVY, True),
        ("range/pair는 split마다 크게 변동", 0, DARKTEXT, False),
        ("(symbol volatility 배치 차이)", 1, GREY, False),
    ],
    size=12,
    gap=7,
)
_box(
    s,
    Inches(0.85),
    Inches(4.95),
    Inches(7.7),
    Inches(1.6),
    fill=RGBColor(0xE7, 0xF3, 0xEC),
    line=GREEN,
    line_w=Pt(1),
)
_text(
    s,
    Inches(1.05),
    Inches(5.08),
    Inches(7.3),
    Inches(1.4),
    [
        [("shape_token은 상대적 candle shape를 안정적으로 포착하고,", 13, NAVY, True)],
        [
            (
                "range_bucket은 symbol별 volatility profile 차이를 민감하게 드러냄",
                13,
                NAVY,
                True,
            )
        ],
        [("→ 'shape + range 분리' 설계를 지지하는 증거", 12, GREEN, True)],
    ],
    space_after=4,
)
footer(s, n)

# ========================================================================
# 12. VOLUME-FILTERED RERUN — 35 runs
# ========================================================================
s = new()
header(s, "PHASE 1B", "Volume-filtered 재실행 — split family별 35 runs")
_text(
    s,
    Inches(0.85),
    Inches(1.7),
    Inches(11.6),
    Inches(0.4),
    [
        [
            (
                "volume ≤ 1 candle 제거 후 random 20 + vol_strat 10 + vol_holdout 5 = 35 runs 재실행",
                13,
                GREY,
                False,
            )
        ]
    ],
)
table(
    s,
    Inches(0.85),
    Inches(2.25),
    Inches(11.6),
    [
        [
            "Metric (test–train L1 mean)",
            "random (20)",
            "vol_strat (10)",
            "vol_holdout (5)",
        ],
        ["shape", "0.085", "0.081", "0.074"],
        ["range", "0.251", "0.102", "0.920"],
        ["pair", "0.289", "0.156", "0.940"],
    ],
    col_w=[3.0, 1.8, 1.8, 1.8],
    font=13,
    highlight_rows={1},
)
bullets(
    s,
    Inches(0.85),
    Inches(4.0),
    Inches(11.6),
    Inches(2.6),
    [
        (
            "shape drift는 세 split family 모두에서 낮게 유지 (0.074 ~ 0.085)",
            0,
            NAVY,
            True,
        ),
        (
            "vol_strat: volatility를 균형 배치하면 range/pair drift도 낮아짐 → random의 drift 상당부분은 shape 문제가 아니라 volatility 구성 차이",
            0,
            DARKTEXT,
            False,
        ),
        (
            "vol_holdout: high-volatility symbol을 train에서 제외 → shape는 안정, range/pair만 ~0.92로 의도대로 급증",
            0,
            DARKTEXT,
            False,
        ),
        ("shape와 volatility context가 분리되고 있다는 해석을 강화", 0, GREEN, True),
    ],
    size=13,
    gap=9,
)
footer(s, n)

# ========================================================================
# 13. VOLUME FILTER EFFECT
# ========================================================================
s = new()
header(s, "ROBUSTNESS", "Volume filter 효과 — 결론은 불변")
table(
    s,
    Inches(0.85),
    Inches(1.85),
    Inches(8.3),
    [
        ["Split / Metric", "Unfiltered", "Vol-filtered", "Δ"],
        ["random · shape", "0.091", "0.085", "−0.006"],
        ["random · pair", "0.292", "0.289", "−0.004"],
        ["vol_strat · shape", "0.096", "0.081", "−0.015"],
        ["vol_strat · pair", "0.164", "0.156", "−0.007"],
        ["vol_holdout · shape", "0.078", "0.074", "−0.004"],
        ["vol_holdout · range", "0.918", "0.920", "+0.001"],
    ],
    col_w=[2.8, 1.5, 1.5, 1.2],
    font=12,
)
bullets(
    s,
    Inches(9.35),
    Inches(2.0),
    Inches(3.2),
    Inches(4),
    [
        ("shape drift 전반적으로 소폭 ↓", 0, NAVY, True),
        ("range/pair drift는 거의 동일", 0, DARKTEXT, False),
        ("기존 결과가 저거래량 candle artifact가 아님을 확인", 0, GREEN, True),
        ("→ volume ≥ 2를 Phase 1/2 공통 전처리로 채택", 0, DARKTEXT, False),
    ],
    size=13,
    gap=11,
)
footer(s, n)

# ========================================================================
# 14. PHASE 2 ENTRY CRITERIA
# ========================================================================
s = new(dark=True)
header(s, "GATE", "Phase 2 진입 기준 검증 (volume-filtered random 20회)", dark=True)
table(
    s,
    Inches(2.4),
    Inches(2.1),
    Inches(8.5),
    [
        ["Criterion", "기준", "결과", "Pass"],
        ["shape test–train L1 mean", "< 0.15", "0.085", "✅"],
        ["shape test–train L1 std", "< 0.05", "0.039", "✅"],
        ["shape test–train L1 max", "< 0.30", "0.157", "✅"],
    ],
    col_w=[3.4, 1.6, 1.6, 1.2],
    font=15,
    header_fill=BLUE,
)
_box(
    s,
    Inches(2.4),
    Inches(4.4),
    Inches(8.5),
    Inches(1.4),
    fill=RGBColor(0x21, 0x3A, 0x2C),
    line=GREEN,
    line_w=Pt(1.25),
)
_text(
    s,
    Inches(2.4),
    Inches(4.62),
    Inches(8.5),
    Inches(1.1),
    [
        [("Phase 1B는 volume-filtered 기준에서도", 17, WHITE, True)],
        [
            (
                "Phase 2 진입 조건을 모두 통과합니다.",
                19,
                RGBColor(0x8C, 0xF0, 0xB8),
                True,
            )
        ],
    ],
    align=PP_ALIGN.CENTER,
    space_after=4,
)
footer(s, n, dark=True)

# ========================================================================
# 15. KEY VISUAL — random split
# ========================================================================
s = new(dark=True)
header(s, "VISUAL EVIDENCE", "대표 random split (random_00)", dark=True)
picture(
    s,
    P1B_RND / "03_shape_range_pair_heatmap.png",
    Inches(0.85),
    Inches(1.85),
    w=Inches(5.7),
    caption="Shape × Range pair heatmap",
)
picture(
    s,
    P1B_RND / "02_range_bucket_ratio_histogram.png",
    Inches(6.9),
    Inches(1.85),
    w=Inches(5.7),
    caption="Range bucket ratio — train vs val vs test",
)
footer(s, n, dark=True)

# ========================================================================
# 16. CONCLUSIONS
# ========================================================================
s = new()
header(s, "CONCLUSION", "결론 및 채택 결정")
_text(
    s,
    Inches(0.85),
    Inches(1.7),
    Inches(5.7),
    Inches(0.4),
    [[("✓ 확인된 것", 15, GREEN, True)]],
)
bullets(
    s,
    Inches(0.85),
    Inches(2.2),
    Inches(5.7),
    Inches(4),
    [
        (
            "price-shape only VQ-VAE token은 여러 random symbol split에서 안정적 shape vocabulary 생성",
            0,
            DARKTEXT,
            False,
        ),
        (
            "range를 별도 bucket으로 분리하면 volatility drift를 shape drift와 구분해 관찰 가능",
            0,
            DARKTEXT,
            False,
        ),
        (
            "shape test drift 평균 0.085 (vol-filtered random), Phase 2 기준 통과",
            0,
            NAVY,
            True,
        ),
        ("volume ≥ 2 filter는 결론을 바꾸지 않음", 0, DARKTEXT, False),
    ],
    size=13,
    gap=10,
)
_text(
    s,
    Inches(6.8),
    Inches(1.7),
    Inches(5.7),
    Inches(0.4),
    [[("⏸ 아직 결론짓지 않은 것", 15, UP, True)]],
)
bullets(
    s,
    Inches(6.8),
    Inches(2.2),
    Inches(5.7),
    Inches(2.6),
    [
        ("shape token이 시장 상태를 설명한다", 0, DARKTEXT, False),
        ("shape token이 미래 dynamics를 예측한다", 0, DARKTEXT, False),
        ("VQ-VAE가 KMeans보다 우월하다", 0, DARKTEXT, False),
    ],
    size=13,
    gap=10,
)
_box(
    s,
    Inches(6.8),
    Inches(4.6),
    Inches(5.7),
    Inches(1.9),
    fill=RGBColor(0xE2, 0xEC, 0xFF),
    line=BLUE,
    line_w=Pt(1),
)
_text(
    s,
    Inches(7.0),
    Inches(4.74),
    Inches(5.3),
    Inches(1.7),
    [
        [("채택 표현", 13, BLUE, True)],
        [("final rep = (shape_token, range_bucket)", 15, NAVY, True)],
        [("codebook K=12 · volume ≥ 2 · KMeans는 baseline 유지", 12, GREY, False)],
    ],
    space_after=5,
)
footer(s, n)

# ========================================================================
# 17. NEXT STEPS
# ========================================================================
s = new(dark=True)
header(s, "NEXT", "다음 단계 — Phase 2 Sequential Dynamics", dark=True)
bullets(
    s,
    Inches(0.85),
    Inches(1.95),
    Inches(11.6),
    Inches(4.5),
    [
        (
            "Phase 1B 표현 (shape_token, range_bucket)을 그대로 Phase 2로 이관",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
        (
            "질문: token sequence가 시간축에서 non-random transition 구조를 갖는가?",
            0,
            WHITE,
            True,
        ),
        (
            "평가: transition matrix · entropy · self-transition · marginal/Markov/shuffled baseline 비교",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
        (
            "공통 전처리 고정: volume ≥ 2,  representation = (shape_token, range_bucket)",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
        (
            "KMeans token vs VQ-VAE token의 transition 구조 비교 지속",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
        (
            "symbol 경계를 넘는 transition 생성 금지 등 leakage 통제 유지",
            0,
            RGBColor(0xD9, 0xE2, 0xF2),
            False,
        ),
    ],
    size=14.5,
    gap=13,
)
_box(s, Inches(0.85), Inches(6.15), Inches(11.6), Inches(0.6), fill=NAVY_2)
_text(
    s,
    Inches(1.05),
    Inches(6.22),
    Inches(11.2),
    Inches(0.5),
    [
        [
            (
                "계획 문서: research/notebooks/02_sequential_dynamics/README.md",
                12,
                DOWN,
                True,
            )
        ]
    ],
)
footer(s, n, dark=True)

prs.save(str(OUT))
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
